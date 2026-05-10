"""PARITY-01 PPTX / XLSX / HTML / MD extractor tests.

All fixtures are programmatic (no checked-in binaries). Each format
covers: text non-empty, image-count expectation, malformed input
raises ValueError.
"""
from __future__ import annotations

import base64
import io
from pathlib import Path

import pytest


def _make_small_png(rgb: tuple[int, int, int], size: int = 80) -> bytes:
    from PIL import Image

    img = Image.new("RGB", (size, size), color=rgb)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_pptx_with_images(tmp_path: Path) -> tuple[Path, int]:
    """3-slide deck, each slide has one picture."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]  # Blank layout

    image_count = 0
    for i in range(3):
        slide = prs.slides.add_slide(blank)
        # Title-ish text via a textbox
        tx = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        tx.text_frame.text = f"Slide {i}: pineapple-bridge {i}"
        # Picture
        png = _make_small_png((255 - i * 50, 0, i * 50), size=120)
        slide.shapes.add_picture(
            io.BytesIO(png),
            Inches(1),
            Inches(2),
            width=Inches(2),
            height=Inches(2),
        )
        image_count += 1

    out = tmp_path / "sample.pptx"
    prs.save(str(out))
    return out, image_count


def test_extract_pptx_text_and_images(
    tmp_project_dir: Path, sample_pptx_with_images
) -> None:
    from nyrahost.extractors.pptx import extract_pptx

    path, count = sample_pptx_with_images
    text, refs = extract_pptx(path, project_saved=tmp_project_dir / "Saved")
    assert "pineapple-bridge" in text
    assert "Slide 0" in text and "Slide 2" in text
    assert len(refs) == count
    for ref in refs:
        assert ref.kind == "image"
        assert Path(ref.path).exists()


def test_extract_pptx_malformed_raises(
    tmp_path: Path, tmp_project_dir: Path
) -> None:
    from nyrahost.extractors.pptx import extract_pptx

    bad = tmp_path / "garbage.pptx"
    bad.write_bytes(b"plain text, not zip")
    with pytest.raises(ValueError):
        extract_pptx(bad, project_saved=tmp_project_dir / "Saved")


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_xlsx(tmp_path: Path) -> tuple[Path, int]:
    """Two-sheet workbook with text data; XLSX rarely has meaningful images."""
    from openpyxl import Workbook

    wb = Workbook()
    ws1 = wb.active
    ws1.title = "Stats"
    ws1.append(["Hero", "DPS", "HP"])
    ws1.append(["Knight", 100, 250])
    ws1.append(["Mage", 180, 120])
    ws1.append(["Rogue", 140, 150])

    ws2 = wb.create_sheet("Notes")
    ws2.append(["pineapple-bridge"])
    ws2.append([None, "second-row-cell"])
    out = tmp_path / "sample.xlsx"
    wb.save(str(out))
    return out, 0  # no embedded images


def test_extract_xlsx_text(tmp_project_dir: Path, sample_xlsx) -> None:
    from nyrahost.extractors.xlsx import extract_xlsx

    path, expected_image_count = sample_xlsx
    text, refs = extract_xlsx(path, project_saved=tmp_project_dir / "Saved")
    assert "pineapple-bridge" in text
    assert "Knight" in text
    assert "DPS" in text
    # Sheet headers are emitted as "# Sheet: <name>".
    assert "# Sheet: Stats" in text
    assert "# Sheet: Notes" in text
    assert len(refs) == expected_image_count


def test_extract_xlsx_malformed_raises(
    tmp_path: Path, tmp_project_dir: Path
) -> None:
    from nyrahost.extractors.xlsx import extract_xlsx

    bad = tmp_path / "garbage.xlsx"
    bad.write_bytes(b"plain text, not zip")
    with pytest.raises(ValueError):
        extract_xlsx(bad, project_saved=tmp_project_dir / "Saved")


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------


def _data_uri_png(rgb: tuple[int, int, int]) -> str:
    png = _make_small_png(rgb, size=80)
    b64 = base64.b64encode(png).decode("ascii")
    return f"data:image/png;base64,{b64}"


@pytest.fixture
def sample_html(tmp_path: Path) -> tuple[Path, int]:
    """Small HTML document with one data-URI image."""
    data_uri = _data_uri_png((10, 200, 10))
    html = (
        "<html><head><title>NyraHTML</title></head>"
        "<body>"
        "<h1>Document Title</h1>"
        "<p>The phrase to find: pineapple-bridge.</p>"
        f'<img src="{data_uri}" alt="green square" />'
        '<img src="https://example.com/skipped.png" alt="not data uri" />'
        "</body></html>"
    )
    out = tmp_path / "sample.html"
    out.write_text(html, encoding="utf-8")
    return out, 1


def test_extract_html_text_and_image(
    tmp_project_dir: Path, sample_html
) -> None:
    from nyrahost.extractors.html import extract_html

    path, expected = sample_html
    text, refs = extract_html(path, project_saved=tmp_project_dir / "Saved")
    assert "pineapple-bridge" in text
    assert "Document Title" in text
    assert len(refs) == expected
    for ref in refs:
        assert ref.kind == "image"
        assert Path(ref.path).exists()


def test_extract_html_empty_raises(tmp_path: Path, tmp_project_dir: Path) -> None:
    from nyrahost.extractors.html import extract_html

    empty = tmp_path / "empty.html"
    empty.write_bytes(b"")
    with pytest.raises(ValueError):
        extract_html(empty, project_saved=tmp_project_dir / "Saved")


# ---------------------------------------------------------------------------
# Markdown
# ---------------------------------------------------------------------------


@pytest.fixture
def sample_md(tmp_path: Path) -> tuple[Path, int]:
    data_uri = _data_uri_png((200, 200, 0))
    md = (
        "# Heading\n\n"
        "Some prose with the phrase pineapple-bridge.\n\n"
        f"![alt text]({data_uri})\n\n"
        "More text after the image.\n"
    )
    out = tmp_path / "sample.md"
    out.write_text(md, encoding="utf-8")
    return out, 1


def test_extract_md_text_and_image(tmp_project_dir: Path, sample_md) -> None:
    from nyrahost.extractors.md import extract_md

    path, expected = sample_md
    text, refs = extract_md(path, project_saved=tmp_project_dir / "Saved")
    assert "pineapple-bridge" in text
    assert "Heading" in text
    assert len(refs) == expected
    for ref in refs:
        assert ref.kind == "image"
        assert Path(ref.path).exists()


def test_extract_md_empty_raises(tmp_path: Path, tmp_project_dir: Path) -> None:
    from nyrahost.extractors.md import extract_md

    empty = tmp_path / "empty.md"
    empty.write_bytes(b"")
    with pytest.raises(ValueError):
        extract_md(empty, project_saved=tmp_project_dir / "Saved")


# ---------------------------------------------------------------------------
# Dispatch — sanity round-trip via the public entrypoint
# ---------------------------------------------------------------------------


def test_dispatch_picks_correct_extractor(
    tmp_path: Path, tmp_project_dir: Path, sample_html
) -> None:
    from nyrahost.extractors import dispatch

    html_path, _ = sample_html
    text, refs = dispatch(html_path, project_saved=tmp_project_dir / "Saved")
    assert "pineapple-bridge" in text
    assert len(refs) >= 1


def test_dispatch_unknown_suffix_raises(
    tmp_path: Path, tmp_project_dir: Path
) -> None:
    from nyrahost.extractors import dispatch

    weird = tmp_path / "anything.unknown"
    weird.write_bytes(b"...")
    with pytest.raises(ValueError):
        dispatch(weird, project_saved=tmp_project_dir / "Saved")
