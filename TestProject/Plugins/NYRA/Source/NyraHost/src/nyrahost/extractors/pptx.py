"""PPTX extractor — python-pptx 1.0.2 (LOCKED-06 approved).

Returns ``(text, list[AttachmentRef])``. Text walks slides → shapes,
collecting ``shape.text_frame`` text; pictures (``shape.shape_type ==
MSO_SHAPE_TYPE.PICTURE``) are extracted via ``shape.image.blob`` and
re-ingested through Pillow normalisation.
"""
from __future__ import annotations

from pathlib import Path
from typing import Tuple

from nyrahost.attachments import AttachmentRef
from nyrahost.extractors._common import (
    assert_ooxml_within_zip_budget,
    ingest_image_bytes,
)


def extract_pptx(
    src: Path, *, project_saved: Path
) -> Tuple[str, list[AttachmentRef]]:
    assert_ooxml_within_zip_budget(src)

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.exc import PackageNotFoundError
    except ImportError as e:
        raise ValueError(
            "python-pptx not installed — cannot extract PPTX documents"
        ) from e

    try:
        prs = Presentation(str(src))
    except PackageNotFoundError as e:
        raise ValueError(f"malformed PPTX {src.name}: {e}") from e
    except Exception as e:
        raise ValueError(f"unreadable PPTX {src.name}: {e}") from e

    text_chunks: list[str] = []
    image_refs: list[AttachmentRef] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line:
                        text_chunks.append(line)
            # Pictures
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                except (AttributeError, ValueError):
                    continue
                if not blob:
                    continue
                ref = ingest_image_bytes(blob, project_saved=project_saved)
                if ref is not None:
                    image_refs.append(ref)
            # Shape groups can also contain pictures (one nesting level
            # is the common case in design decks); recurse one level.
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for inner in shape.shapes:
                    if inner.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            blob = inner.image.blob
                        except (AttributeError, ValueError):
                            continue
                        if not blob:
                            continue
                        ref = ingest_image_bytes(
                            blob, project_saved=project_saved
                        )
                        if ref is not None:
                            image_refs.append(ref)

    text = "\n".join(text_chunks)
    return text, image_refs
